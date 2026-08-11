#!/usr/bin/env python3
"""Generate a one-slide, editable research summary for WAM alignment."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("WAM_video_action_alignment_summary_zh.pptx")

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

COLORS = {
    "bg": "F5F8FC",
    "paper": "FFFFFF",
    "navy": "11253E",
    "ink": "1A2B3C",
    "muted": "65758B",
    "line": "DCE4ED",
    "wo": "728093",
    "wo_bg": "EDF1F5",
    "w": "0E9895",
    "w_bg": "E4F7F5",
    "positive": "15986A",
    "positive_bg": "E6F6EF",
    "negative": "D9564E",
    "negative_bg": "FCECEA",
    "purple": "6750A4",
    "purple_bg": "F0ECFA",
    "blue": "3767C7",
    "blue_bg": "EAF0FC",
    "amber": "BD7B16",
    "amber_bg": "FFF3DA",
    "white_75": "CAD4DF",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_shape_name(shape, name: str) -> None:
    shape._element.nvSpPr.cNvPr.set("name", name)


def set_run_font(run, size, color, bold=False, font_name=FONT_CN):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")
    for tag, face in (("latin", FONT_EN), ("ea", FONT_CN), ("cs", FONT_EN)):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = OxmlElement(f"a:{tag}")
            rpr.append(node)
        node.set("typeface", face)


def text_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    *,
    size=12,
    color="ink",
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0,
    font_name=FONT_CN,
    name=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        set_shape_name(box, name)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    set_run_font(run, size, COLORS[color], bold=bold, font_name=font_name)
    return box


def rich_text_box(
    slide,
    x,
    y,
    w,
    h,
    spans,
    *,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0,
    name=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        set_shape_name(box, name)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        set_run_font(
            run,
            span.get("size", 12),
            COLORS[span.get("color", "ink")],
            bold=span.get("bold", False),
            font_name=span.get("font_name", FONT_CN),
        )
    return box


def rounded_rect(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill="paper",
    line=None,
    radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    name=None,
):
    shape = slide.shapes.add_shape(
        radius_shape, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        set_shape_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    if line:
        shape.line.color.rgb = rgb(COLORS[line])
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def pill(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    fill,
    color,
    size=9,
    bold=True,
    line=None,
    name=None,
):
    rounded_rect(
        slide, x, y, w, h, fill=fill, line=line, name=f"{name}_shape" if name else None
    )
    return text_box(
        slide,
        x,
        y,
        w,
        h,
        text,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        name=f"{name}_text" if name else None,
    )


def tiny_label(slide, x, y, w, text, color="muted", name=None):
    return text_box(
        slide,
        x,
        y,
        w,
        0.18,
        text,
        size=7.5,
        color=color,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        name=name,
    )


def divider(slide, x1, y1, x2, y2, color="line", width=0.8, name=None):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1),
        Inches(y1),
        Inches(max(x2 - x1, 0.006)),
        Inches(max(y2 - y1, 0.006)),
    )
    if name:
        set_shape_name(line, name)
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS[color])
    line.line.fill.background()
    return line


def comparison_row(
    slide,
    x,
    y,
    label,
    wo_value,
    w_value,
    delta,
    *,
    positive,
    name,
):
    text_box(
        slide,
        x,
        y,
        0.92,
        0.34,
        label,
        size=10.2,
        color="ink",
        bold=True,
        name=f"{name}_label",
    )
    pill(
        slide,
        x + 1.02,
        y,
        0.86,
        0.34,
        wo_value,
        fill="wo_bg",
        color="wo",
        size=10,
        name=f"{name}_wo",
    )
    text_box(
        slide,
        x + 1.91,
        y,
        0.30,
        0.34,
        "→",
        size=12,
        color="muted",
        bold=True,
        align=PP_ALIGN.CENTER,
        name=f"{name}_arrow",
    )
    pill(
        slide,
        x + 2.25,
        y,
        0.86,
        0.34,
        w_value,
        fill="w_bg",
        color="w",
        size=10,
        name=f"{name}_w",
    )
    pill(
        slide,
        x + 3.25,
        y,
        0.94,
        0.34,
        delta,
        fill="positive_bg" if positive else "negative_bg",
        color="positive" if positive else "negative",
        size=9.5,
        name=f"{name}_delta",
    )


def evidence_card(
    slide,
    x,
    y,
    title,
    subtitle,
    best_text,
    row1,
    row2,
    footnote,
    *,
    positive,
    name,
):
    rounded_rect(slide, x, y, 5.64, 1.84, fill="paper", name=f"{name}_card")
    text_box(
        slide,
        x + 0.18,
        y + 0.12,
        2.00,
        0.25,
        title,
        size=13,
        color="navy",
        bold=True,
        name=f"{name}_title",
    )
    text_box(
        slide,
        x + 1.55,
        y + 0.14,
        1.74,
        0.20,
        subtitle,
        size=7.5,
        color="muted",
        name=f"{name}_subtitle",
    )
    pill(
        slide,
        x + 3.52,
        y + 0.10,
        1.94,
        0.32,
        best_text,
        fill="positive_bg" if positive else "blue_bg",
        color="positive" if positive else "blue",
        size=8.2,
        name=f"{name}_best",
    )

    tiny_label(slide, x + 1.22, y + 0.45, 0.55, "wo", "wo")
    tiny_label(slide, x + 2.45, y + 0.45, 0.55, "w", "w")
    tiny_label(slide, x + 3.45, y + 0.45, 0.65, "Δ", "muted")
    comparison_row(
        slide,
        x + 0.18,
        y + 0.62,
        row1[0],
        row1[1],
        row1[2],
        row1[3],
        positive=positive,
        name=f"{name}_row1",
    )
    comparison_row(
        slide,
        x + 0.18,
        y + 1.04,
        row2[0],
        row2[1],
        row2[2],
        row2[3],
        positive=positive,
        name=f"{name}_row2",
    )
    text_box(
        slide,
        x + 0.18,
        y + 1.48,
        5.28,
        0.22,
        footnote,
        size=8,
        color="muted",
        name=f"{name}_note",
    )


def insight_row(slide, x, y, number, text, accent, name):
    pill(
        slide,
        x,
        y,
        0.31,
        0.31,
        number,
        fill=accent,
        color="paper",
        size=8.5,
        name=f"{name}_number",
    )
    text_box(
        slide,
        x + 0.43,
        y - 0.02,
        3.17,
        0.43,
        text,
        size=9.6,
        color="paper",
        bold=False,
        valign=MSO_ANCHOR.TOP,
        name=f"{name}_text",
    )


def roadmap_stage(slide, x, y, w, number, title, body, accent, name):
    pill(
        slide,
        x,
        y + 0.03,
        0.32,
        0.32,
        number,
        fill=accent,
        color="paper",
        size=8.5,
        name=f"{name}_step",
    )
    text_box(
        slide,
        x + 0.43,
        y,
        w - 0.43,
        0.25,
        title,
        size=10.4,
        color="navy",
        bold=True,
        name=f"{name}_title",
    )
    text_box(
        slide,
        x + 0.43,
        y + 0.28,
        w - 0.43,
        0.38,
        body,
        size=8.4,
        color="muted",
        valign=MSO_ANCHOR.TOP,
        name=f"{name}_body",
    )


def build_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = (
        "从固定耦合到自适应对齐：WAM Video–Action Alignment 的初步证据与下一步"
    )
    prs.core_properties.subject = "LIBERO and RoboTwin 2.0 incremental alignment study"
    prs.core_properties.author = "FAST_WAM"
    prs.core_properties.keywords = "WAM, video-action alignment, LIBERO, RoboTwin, REFL, gate"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])

    # Header
    divider(slide, 0.35, 0.25, 0.39, 0.95, color="w", width=1.0, name="accent_bar")
    text_box(
        slide,
        0.55,
        0.22,
        4.50,
        0.18,
        "WAM  •  VIDEO–ACTION ALIGNMENT",
        size=8,
        color="w",
        bold=True,
        font_name=FONT_EN,
        name="overline",
    )
    text_box(
        slide,
        0.55,
        0.43,
        10.85,
        0.42,
        "从固定耦合到自适应对齐：初步证据与下一步",
        size=25,
        color="navy",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        name="slide_title",
    )
    text_box(
        slide,
        0.55,
        0.87,
        10.80,
        0.24,
        "核心问题不是“要不要 video”，而是 action 何时、以多大强度相信预测未来",
        size=10.5,
        color="muted",
        name="slide_subtitle",
    )
    pill(
        slide,
        11.47,
        0.30,
        1.43,
        0.34,
        "PRELIMINARY",
        fill="amber_bg",
        color="amber",
        size=8.5,
        name="preliminary_badge",
    )

    # Left: research setup
    rounded_rect(slide, 0.35, 1.25, 2.52, 4.20, fill="paper", name="setup_card")
    text_box(
        slide,
        0.55,
        1.42,
        1.75,
        0.25,
        "01  前期探索",
        size=12,
        color="navy",
        bold=True,
        name="setup_title",
    )
    pill(
        slide,
        2.05,
        1.41,
        0.58,
        0.27,
        "DONE",
        fill="positive_bg",
        color="positive",
        size=7.5,
        name="done_badge",
    )
    text_box(
        slide,
        0.55,
        1.78,
        2.10,
        0.52,
        "ActionDiT 应读取多少\n生成视频信息？",
        size=14,
        color="ink",
        bold=True,
        valign=MSO_ANCHOR.TOP,
        name="research_question",
    )

    tiny_label(slide, 0.55, 2.47, 1.6, "CONTEXT LEVEL", "muted", "context_label")
    pill(
        slide,
        0.55,
        2.75,
        0.50,
        0.35,
        "wo",
        fill="wo_bg",
        color="wo",
        size=10,
        name="context_wo",
    )
    text_box(
        slide,
        1.08,
        2.75,
        0.25,
        0.35,
        "→",
        size=10,
        color="muted",
        align=PP_ALIGN.CENTER,
        name="context_arrow1",
    )
    pill(
        slide,
        1.35,
        2.75,
        0.77,
        0.35,
        "prefix N",
        fill="purple_bg",
        color="purple",
        size=8.5,
        name="context_prefix",
    )
    text_box(
        slide,
        2.14,
        2.75,
        0.22,
        0.35,
        "→",
        size=10,
        color="muted",
        align=PP_ALIGN.CENTER,
        name="context_arrow2",
    )
    pill(
        slide,
        2.37,
        2.75,
        0.32,
        0.35,
        "w",
        fill="w_bg",
        color="w",
        size=10,
        name="context_w",
    )
    text_box(
        slide,
        0.55,
        3.17,
        2.05,
        0.38,
        "first frame only  →  full video tokens",
        size=7.8,
        color="muted",
        font_name=FONT_EN,
        name="context_definition",
    )

    tiny_label(slide, 0.55, 3.66, 1.8, "PARAMETER LEVEL", "muted", "parameter_label")
    pill(
        slide,
        0.55,
        3.94,
        0.92,
        0.36,
        "Shared",
        fill="blue_bg",
        color="blue",
        size=9,
        name="shared_chip",
    )
    pill(
        slide,
        1.55,
        3.94,
        1.10,
        0.36,
        "TwoAction",
        fill="purple_bg",
        color="purple",
        size=9,
        name="twoaction_chip",
    )
    text_box(
        slide,
        0.55,
        4.37,
        2.02,
        0.35,
        "共享参数  vs.  分离 wo / w experts",
        size=8,
        color="muted",
        name="parameter_definition",
    )
    divider(slide, 0.55, 4.79, 2.66, 4.80, color="line", name="setup_divider")
    text_box(
        slide,
        0.55,
        4.91,
        2.05,
        0.33,
        "LIBERO  +  RoboTwin 2.0",
        size=10.2,
        color="navy",
        bold=True,
        name="benchmarks_label",
    )
    text_box(
        slide,
        0.55,
        5.17,
        2.05,
        0.20,
        "固定策略已完成系统性消融",
        size=7.8,
        color="muted",
        name="benchmark_note",
    )

    # Middle: evidence
    text_box(
        slide,
        3.06,
        1.25,
        3.20,
        0.26,
        "02  关键证据  ·  Success Rate",
        size=12,
        color="navy",
        bold=True,
        name="evidence_title",
    )
    evidence_card(
        slide,
        3.06,
        1.61,
        "LIBERO",
        "2,000 episodes / mode",
        "BEST  Shared + w  98.50%",
        ("Shared", "97.85", "98.50", "+0.65pp"),
        ("TwoAction", "97.30", "98.40", "+1.10pp"),
        "LIBERO-10：Shared w 比 wo +2.6pp；prefix 1–9/10 未呈现单调收益",
        positive=True,
        name="libero",
    )
    evidence_card(
        slide,
        3.06,
        3.61,
        "RoboTwin 2.0",
        "clean/random avg · 10k/mode",
        "BEST  Shared + wo  93.15%",
        ("Shared", "93.15", "91.06", "−2.09pp"),
        ("TwoAction", "92.18", "91.57", "−0.61pp"),
        "Shared 的额外失败高度集中：约 84% 来自 2 个任务 → reliability 明显 task-dependent",
        positive=False,
        name="robotwin",
    )

    # Right: central insight
    rounded_rect(slide, 8.88, 1.25, 4.10, 4.20, fill="navy", name="insight_card")
    text_box(
        slide,
        9.13,
        1.43,
        1.65,
        0.20,
        "03  CORE INSIGHT",
        size=8.3,
        color="w",
        bold=True,
        font_name=FONT_EN,
        name="insight_overline",
    )
    rich_text_box(
        slide,
        9.13,
        1.77,
        3.55,
        0.80,
        [
            {"text": "Video 对 action 的价值\n", "size": 17, "color": "paper", "bold": True},
            {
                "text": "取决于 benchmark / task",
                "size": 17,
                "color": "w",
                "bold": True,
            },
        ],
        valign=MSO_ANCHOR.TOP,
        name="insight_headline",
    )
    text_box(
        slide,
        9.13,
        2.64,
        3.56,
        0.43,
        "预测未来既可能是 foresight，\n也可能成为 structured noise。",
        size=9.7,
        color="white_75",
        valign=MSO_ANCHOR.TOP,
        name="insight_thesis",
    )
    divider(slide, 9.13, 3.14, 12.69, 3.15, color="wo", name="insight_divider")
    insight_row(
        slide,
        9.13,
        3.34,
        "1",
        "固定 w 可带来增益，也会引入负迁移",
        "w",
        "insight_1",
    )
    insight_row(
        slide,
        9.13,
        3.89,
        "2",
        "TwoAction 只能缓解，未稳定优于 Shared",
        "purple",
        "insight_2",
    )
    insight_row(
        slide,
        9.13,
        4.44,
        "3",
        "固定 prefix 时长尚未找到稳定 sweet spot",
        "amber",
        "insight_3",
    )
    pill(
        slide,
        9.13,
        5.00,
        3.56,
        0.32,
        "→ reliability-aware adaptive alignment",
        fill="w",
        color="paper",
        size=9,
        name="adaptive_alignment_callout",
    )

    # Bottom: roadmap
    rounded_rect(slide, 0.35, 5.66, 12.63, 1.42, fill="paper", name="roadmap_card")
    text_box(
        slide,
        0.55,
        5.78,
        4.90,
        0.25,
        "下一步  ·  从 static choice 到 adaptive alignment",
        size=12,
        color="navy",
        bold=True,
        name="roadmap_title",
    )
    pill(
        slide,
        8.58,
        5.77,
        4.10,
        0.29,
        "目标：保留 LIBERO 增益，避免 RoboTwin 负迁移",
        fill="w_bg",
        color="w",
        size=8.5,
        name="roadmap_target",
    )
    roadmap_stage(
        slide,
        0.55,
        6.18,
        3.50,
        "1",
        "REFL-style alignment",
        "用任务回报 / consistency 构造 video relevance 信号，快速验证“可对齐性”",
        "purple",
        "roadmap_refl",
    )
    divider(slide, 4.17, 6.15, 4.18, 6.86, color="line", name="roadmap_divider1")
    roadmap_stage(
        slide,
        4.37,
        6.18,
        3.72,
        "2",
        "Learned adaptive gate",
        "学习 g(state, task, t, confidence) ∈ [0,1]，连续混合 wo ↔ w",
        "w",
        "roadmap_gate",
    )
    divider(slide, 8.22, 6.15, 8.23, 6.86, color="line", name="roadmap_divider2")
    roadmap_stage(
        slide,
        8.42,
        6.18,
        4.26,
        "3",
        "验证闭环",
        "per-task oracle → calibration / failure attribution → ≥3 seeds、OOD、cost–success",
        "blue",
        "roadmap_validation",
    )

    # Footnote
    text_box(
        slide,
        0.42,
        7.17,
        12.46,
        0.18,
        "Preliminary · final checkpoint / fixed evaluation · no multi-seed CI · no external SOTA claim  |  Results validated 2026-07",
        size=7.3,
        color="muted",
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
        name="footnote",
    )

    return prs


if __name__ == "__main__":
    deck = build_slide()
    deck.save(OUT)
    print(OUT)
